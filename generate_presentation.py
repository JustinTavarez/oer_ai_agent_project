"""Generate the OER AI Agent project presentation PowerPoint file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DARK_BG = RGBColor(0x1E, 0x1B, 0x2E)
PURPLE_ACCENT = RGBColor(0x7C, 0x3A, 0xED)
LIGHT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xD1, 0xD5, 0xDB)
INDIGO = RGBColor(0x63, 0x66, 0xF1)
TEAL = RGBColor(0x14, 0xB8, 0xA6)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(8)
        p.level = 0
    return tf


def add_accent_bar(slide, left, top, width, height, color=PURPLE_ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_rounded_box(slide, left, top, width, height, fill_color, text="",
                    font_size=14, text_color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.color.rgb = text_color
        tf.paragraphs[0].font.name = "Calibri"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # ── Slide 1: Title ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), PURPLE_ACCENT)
    add_accent_bar(slide, Inches(0), Inches(7.42), SLIDE_WIDTH, Inches(0.08), INDIGO)
    add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
                "OER AI Agent", font_size=54, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8),
                "An AI-Powered Platform for Discovering Open Educational Resources",
                font_size=24, color=LIGHT_PURPLE, alignment=PP_ALIGN.CENTER)
    add_accent_bar(slide, Inches(5.5), Inches(4.3), Inches(2.3), Inches(0.05), LIGHT_PURPLE)
    add_textbox(slide, Inches(2), Inches(4.8), Inches(9), Inches(0.6),
                "Project Presentation", font_size=20, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)

    # ── Slide 2: Why I Chose This Topic ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_HEIGHT, PURPLE_ACCENT)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                "Why I Chose This Topic", font_size=36, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(3), Inches(0.04), LIGHT_PURPLE)

    items = [
        "\u2022  Textbook costs are a major financial barrier for students",
        "\u2022  Open Educational Resources (OER) provide free, openly licensed materials",
        "\u2022  Finding the right OER is difficult with traditional keyword search",
        "\u2022  AI can transform resource discovery into a guided conversation",
        "\u2022  Passionate about making education accessible and affordable",
        "\u2022  Excited to explore real-world AI applications in education",
    ]
    add_bullet_list(slide, Inches(1), Inches(1.6), Inches(10.5), Inches(5),
                    items, font_size=20, color=LIGHT_GRAY)

    # ── Slide 3: Why It's Important ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_HEIGHT, INDIGO)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                "Why It\u2019s Important", font_size=36, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(3), Inches(0.04), LIGHT_PURPLE)

    box_data = [
        ("Cost Reduction", "Students spend hundreds per semester on textbooks.\nOER eliminates this financial burden."),
        ("Accessibility", "Free, openly licensed content that anyone\ncan use, adapt, and redistribute."),
        ("Smarter Discovery", "AI replaces keyword search with natural\nlanguage conversation for better results."),
    ]
    for i, (title, desc) in enumerate(box_data):
        x = Inches(0.8 + i * 4.1)
        add_rounded_box(slide, x, Inches(1.8), Inches(3.6), Inches(1.0),
                        PURPLE_ACCENT, title, font_size=20, text_color=WHITE)
        add_textbox(slide, x + Inches(0.2), Inches(3.0), Inches(3.2), Inches(1.6),
                    desc, font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(1.5),
                "\"Instead of sifting through dozens of search results, a student can simply ask "
                "for what they need in plain language and receive curated, vetted recommendations.\"",
                font_size=17, color=LIGHT_PURPLE, alignment=PP_ALIGN.CENTER)

    # ── Slide 4: What I Want to Learn ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_HEIGHT, TEAL)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                "What I Want to Learn", font_size=36, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(3), Inches(0.04), TEAL)

    learn_items = [
        ("\U0001F916  LLM Integration", "How to connect large language models to real applications using FastAPI and OpenAI-compatible APIs"),
        ("\U0001F3A8  Modern Front-End Development", "Building responsive, animated interfaces with React 19, Vite, and Tailwind CSS"),
        ("\U0001F4DD  Prompt Engineering", "Designing system prompts that produce structured, consistent JSON output from AI models"),
        ("\U0001F527  Full-Stack Architecture", "RESTful API design, CORS management, and deploying a polished end-to-end application"),
    ]
    for i, (title, desc) in enumerate(learn_items):
        y = Inches(1.6 + i * 1.35)
        add_rounded_box(slide, Inches(0.8), y, Inches(4.2), Inches(1.0),
                        RGBColor(0x2D, 0x28, 0x45), title, font_size=18, text_color=TEAL)
        add_textbox(slide, Inches(5.3), y + Inches(0.15), Inches(7.2), Inches(0.8),
                    desc, font_size=16, color=LIGHT_GRAY)

    # ── Slide 5: What I Want to Share ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_HEIGHT, PURPLE_ACCENT)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                "What I Want to Share with the Audience",
                font_size=36, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(5), Inches(0.04), LIGHT_PURPLE)

    share_items = [
        "\u2022  AI can solve meaningful, real-world problems beyond simple chatbots",
        "\u2022  The technical architecture: how React, FastAPI, and LM Studio work together",
        "\u2022  How natural language processing can replace tedious keyword search",
        "\u2022  Creative thinking about using AI to reduce costs and improve access to education",
        "\u2022  Practical lessons in building and deploying a full-stack AI application",
    ]
    add_bullet_list(slide, Inches(1), Inches(1.6), Inches(10.5), Inches(4.5),
                    share_items, font_size=20, color=LIGHT_GRAY)

    # ── Slide 6: Technical Architecture ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_HEIGHT, INDIGO)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                "Technical Architecture", font_size=36, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(3), Inches(0.04), INDIGO)

    layers = [
        ("Frontend", "React 19  |  Vite 7  |  Tailwind CSS 4\nFramer Motion  |  React Router 7", PURPLE_ACCENT),
        ("Backend", "FastAPI  |  Pydantic\nhttpx  |  RESTful API", INDIGO),
        ("AI Layer", "LM Studio  |  OpenAI-Compatible API\nMeta Llama 3.1 8B Instruct", TEAL),
    ]
    for i, (title, desc, color) in enumerate(layers):
        y = Inches(1.6 + i * 1.8)
        add_rounded_box(slide, Inches(1), y, Inches(3), Inches(1.3),
                        color, title, font_size=22, text_color=WHITE)
        add_textbox(slide, Inches(4.5), y + Inches(0.2), Inches(7.5), Inches(1.0),
                    desc, font_size=17, color=LIGHT_GRAY)
        if i < len(layers) - 1:
            add_textbox(slide, Inches(2.2), y + Inches(1.3), Inches(0.6), Inches(0.5),
                        "\u25BC", font_size=20, color=LIGHT_PURPLE, alignment=PP_ALIGN.CENTER)

    # ── Slide 7: Key Features ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_HEIGHT, PURPLE_ACCENT)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                "Key Features", font_size=36, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(3), Inches(0.04), LIGHT_PURPLE)

    features = [
        ("Conversational AI Search", "Ask for resources in plain English\nand get curated recommendations"),
        ("Course & Source Filters", "Filter by specific courses (ENGL 1101,\nHIST 2111, etc.) and OER sources"),
        ("Structured Resource Cards", "Title, match reason, license, quality\nsummary, and instructor ideas"),
        ("Real-Time Status Monitor", "Live connectivity checks for backend\nand LM Studio every 30 seconds"),
        ("Search History", "Persistent local storage of up to\n15 recent searches"),
        ("Modern Dark UI", "Purple/indigo gradient theme with\nsmooth Framer Motion animations"),
    ]
    for i, (title, desc) in enumerate(features):
        col = i % 3
        row = i // 3
        x = Inches(0.6 + col * 4.15)
        y = Inches(1.6 + row * 2.7)
        add_rounded_box(slide, x, y, Inches(3.8), Inches(0.7),
                        PURPLE_ACCENT, title, font_size=17, text_color=WHITE)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.85), Inches(3.5), Inches(1.2),
                    desc, font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ── Slide 8: Live Demo Plan ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_HEIGHT, TEAL)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                "What Will Be Showcased", font_size=36, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(3), Inches(0.04), TEAL)

    steps = [
        "1.  Landing Page walkthrough \u2014 Hero, Features, How It Works sections",
        "2.  Chat Interface demo \u2014 asking natural language questions for OER",
        "3.  Course & source filtering \u2014 narrowing results by course and provider",
        "4.  AI-generated Resource Cards \u2014 structured recommendations with details",
        "5.  System architecture overview \u2014 React \u2192 FastAPI \u2192 LM Studio pipeline",
        "6.  Prompt engineering deep-dive \u2014 how the system prompt produces JSON output",
        "7.  Challenges, lessons learned, and future enhancements",
    ]
    add_bullet_list(slide, Inches(1), Inches(1.6), Inches(10.5), Inches(5),
                    steps, font_size=19, color=LIGHT_GRAY)

    # ── Slide 9: Future Enhancements ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_HEIGHT, INDIGO)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                "Future Enhancements", font_size=36, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(3), Inches(0.04), INDIGO)

    future_items = [
        ("\U0001F5C4\uFE0F  Database-Backed Indexing", "Store and index OER resources for faster, more reliable retrieval"),
        ("\U0001F512  User Authentication", "Allow users to save favorites, track search history across devices"),
        ("\U0001F310  Expanded Sources", "Integrate additional OER repositories beyond GGC Syllabi and Open ALG"),
        ("\U0001F4CA  Analytics Dashboard", "Track popular searches, resource usage, and user engagement metrics"),
    ]
    for i, (title, desc) in enumerate(future_items):
        y = Inches(1.6 + i * 1.3)
        add_rounded_box(slide, Inches(0.8), y, Inches(4.5), Inches(0.9),
                        RGBColor(0x2D, 0x28, 0x45), title, font_size=18, text_color=INDIGO)
        add_textbox(slide, Inches(5.6), y + Inches(0.15), Inches(7), Inches(0.7),
                    desc, font_size=16, color=LIGHT_GRAY)

    # ── Slide 10: Thank You ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), PURPLE_ACCENT)
    add_accent_bar(slide, Inches(0), Inches(7.42), SLIDE_WIDTH, Inches(0.08), INDIGO)
    add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.2),
                "Thank You!", font_size=54, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2), Inches(3.6), Inches(9), Inches(0.8),
                "Questions & Discussion", font_size=28, color=LIGHT_PURPLE,
                alignment=PP_ALIGN.CENTER)
    add_accent_bar(slide, Inches(5.5), Inches(4.6), Inches(2.3), Inches(0.05), LIGHT_PURPLE)
    add_textbox(slide, Inches(2), Inches(5.0), Inches(9), Inches(0.6),
                "OER AI Agent \u2014 Making Education Accessible Through AI",
                font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    prs.save("/workspace/OER_AI_Agent_Presentation.pptx")
    print("Presentation saved to /workspace/OER_AI_Agent_Presentation.pptx")


if __name__ == "__main__":
    build_presentation()
